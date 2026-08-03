#!/usr/bin/env python3
"""Build UPPMS graduation presentation (English) as PPTX."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ASSETS = Path(__file__).resolve().parent
OUT = ASSETS / "UPPMS-Presentation.pptx"

W, H = Inches(13.333), Inches(7.5)

NAVY = RGBColor(0x0B, 0x2A, 0x3F)
TEAL = RGBColor(0x0D, 0x7A, 0x6F)
ACCENT = RGBColor(0x1F, 0xA3, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF7, 0xF8)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x4A, 0x5A, 0x64)
ORANGE = RGBColor(0xB5, 0x4A, 0x24)
CARD_LINE = RGBColor(0xD0, 0xDC, 0xE0)


def set_run(run, text, size=20, bold=False, color=DARK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def add_top_bar(slide, title: str):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    set_run(r, "  " + title, 26, True, WHITE)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.95), W, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()


def add_footer(slide, page: int, total: int):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(10), Inches(0.35))
    r = box.text_frame.paragraphs[0].add_run()
    set_run(r, "UPPMS  ·  Graduation Project Presentation", 11, False, MUTED)
    num = slide.shapes.add_textbox(Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.35))
    p2 = num.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    set_run(r2, f"{page} / {total}", 11, False, MUTED)


def add_bullets(slide, left, top, width, height, items, size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run()
        set_run(r, "•  " + item, size, False, DARK)


def add_card(slide, left, top, width, height, title, body_lines, title_color=TEAL):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = CARD_LINE
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, title, 18, True, title_color)
    for line in body_lines:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        r2 = p2.add_run()
        set_run(r2, line, 14, False, DARK)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_image_fit(slide, path: Path, left, top, max_w, max_h):
    if not path.exists():
        box = slide.shapes.add_textbox(left, top, max_w, Inches(0.5))
        r = box.text_frame.paragraphs[0].add_run()
        set_run(r, f"[Missing: {path.name}]", 14, False, MUTED)
        return None
    pic = slide.shapes.add_picture(str(path), left, top, width=max_w)
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = max_h
        pic.left = int(left + (max_w - pic.width) / 2)
    return pic


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    gantt = ASSETS / "gantt Chart New.png"
    activity = ASSETS / "activity.png"
    sequence = ASSETS / "sequence.png"
    usecase = ASSETS / "use-case-overview.png"
    arch = ASSETS / "02-architecture-layers.png"
    erd = ASSETS / "03-erd-overview.png"

    # 1 Title
    s = blank_slide(prs)
    add_bg(s, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.6), W, Inches(0.9))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    title = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.0))
    r = title.text_frame.paragraphs[0].add_run()
    set_run(r, "UPPMS", 54, True, WHITE)

    sub = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(0.6))
    r = sub.text_frame.paragraphs[0].add_run()
    set_run(r, "University Project Portfolio Management System", 26, False, ACCENT)

    names = s.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.5), Inches(2.6))
    tf = names.text_frame
    tf.word_wrap = True
    lines = [
        ("Students", True, WHITE, 14),
        ("Baraa Abdullah", False, WHITE, 22),
        ("Sohayb Sabboura", False, WHITE, 22),
        ("Obada Kahwajy", False, WHITE, 22),
        ("", False, WHITE, 8),
        ("Supervised by", True, ACCENT, 14),
        ("Anas Abdulaziz", False, WHITE, 22),
    ]
    first = True
    for text, bold, color, size in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        set_run(r, text, size, bold, color)
        p.space_after = Pt(2)

    foot = s.shapes.add_textbox(Inches(0.8), Inches(6.75), Inches(11), Inches(0.4))
    r = foot.text_frame.paragraphs[0].add_run()
    set_run(r, "Graduation Project Presentation  ·  Syrian Private University", 14, False, WHITE)

    # 2 Agenda
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "Agenda")
    add_bullets(
        s,
        Inches(0.9),
        Inches(1.4),
        Inches(11),
        Inches(5),
        [
            "1.  Problem Statement",
            "2.  How Our System Solves the Problem",
            "3.  System Advantages",
            "4.  What Was Implemented Before vs Now",
            "5.  Methodology (Gantt, Use Case, Activity, Sequence)",
            "6.  Tools Used",
            "7.  Live Demo",
        ],
        22,
    )

    # 3 Problem
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "1. Problem Statement")
    add_bullets(
        s,
        Inches(0.8),
        Inches(1.3),
        Inches(11.5),
        Inches(5.5),
        [
            "Graduation projects are managed with scattered tools (chat, sheets, email).",
            "No reliable isolation when multiple universities share one platform.",
            "Open registration allows unverified accounts into the academic workspace.",
            "Choosing a project idea and breaking it into tasks is slow and unstructured.",
            "Defense scheduling is manual, conflict-prone, and hard to optimize.",
            "Academic progression across stages lacks enforced prerequisites.",
            "Supervisors and admins lack a single source of truth for proposals, progress, and defenses.",
        ],
        18,
    )

    # 4 Solution
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "2. How Our System Solves the Problem")
    add_card(
        s,
        Inches(0.5),
        Inches(1.3),
        Inches(4.0),
        Inches(2.4),
        "Secure Access",
        [
            "Multi-tenant isolation by university",
            "XML-based authorized registration",
            "Role-based access (RBAC)",
        ],
    )
    add_card(
        s,
        Inches(4.7),
        Inches(1.3),
        Inches(4.0),
        Inches(2.4),
        "Smart Project Flow",
        [
            "AI project ideation (Gemini)",
            "Supervisor proposal approval",
            "AI task breakdown + GitHub link",
        ],
    )
    add_card(
        s,
        Inches(8.9),
        Inches(1.3),
        Inches(3.9),
        Inches(2.4),
        "Academic Control",
        [
            "Configurable academic tracks",
            "Prerequisite stage unlocking",
            "Defense result updates progress",
        ],
    )
    add_card(
        s,
        Inches(0.5),
        Inches(4.0),
        Inches(6.1),
        Inches(2.4),
        "Intelligent Defense Scheduling",
        [
            "Genetic algorithm generates conflict-aware schedules",
            "Rooms, availability, committees, and constraints",
            "Admin reviews candidates and approves the final plan",
        ],
    )
    add_card(
        s,
        Inches(6.9),
        Inches(4.0),
        Inches(5.9),
        Inches(2.4),
        "One Platform",
        [
            "Students, supervisors, university & platform admins",
            "Projects, tasks, notifications, dashboards together",
            "Spec-driven delivery for reliable increments",
        ],
    )

    # 5 Advantages
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "3. System Advantages")
    add_bullets(
        s,
        Inches(0.8),
        Inches(1.3),
        Inches(11.5),
        Inches(5.5),
        [
            "University-level data isolation with shared SaaS deployment",
            "Verified onboarding via official XML records (student: email + number, supervisor: email)",
            "End-to-end lifecycle: idea → proposal → project → tasks → defense",
            "AI assistance for ideation and task decomposition",
            "Optimized defense scheduling with hard and soft constraints",
            "Manual committees when needed; genetic search when scale is high",
            "Academic tracks enforce sequential progress and prevent skipping stages",
            "Modern React SPA + Laravel API — maintainable and extensible",
        ],
        18,
    )

    # 6 Before vs Now
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "4. Implemented Before vs Now")
    add_card(
        s,
        Inches(0.5),
        Inches(1.3),
        Inches(6.0),
        Inches(5.3),
        "Previously Implemented",
        [
            "• Multi-tenancy foundation (universities)",
            "• RBAC & account approval workflow",
            "• UI workspace isolation",
            "• Core projects, tasks, comments",
            "• Invitations & GitHub linking",
            "• Notifications & dashboards",
            "• AI project ideation",
            "• AI task breakdown + flexible UI",
        ],
        TEAL,
    )
    add_card(
        s,
        Inches(6.8),
        Inches(1.3),
        Inches(6.0),
        Inches(5.3),
        "Implemented Now (Recent Phases)",
        [
            "• Scheduling infrastructure (rooms, availability, stages)",
            "• Genetic Algorithm scheduling engine",
            "• Scheduling API & admin dashboard",
            "• XML registration & data matching",
            "• Project proposals & supervisor approval",
            "• Manual defense committee management",
            "• Academic tracks & sequential progress",
            "• Full documentation & design diagrams",
        ],
        ORANGE,
    )

    # 7 Methodology
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "5. Methodology")
    add_bullets(
        s,
        Inches(0.8),
        Inches(1.3),
        Inches(11.5),
        Inches(5.5),
        [
            "Spec-Driven Development (Spec Kit): constitution → specify → plan → tasks → implement",
            "Incremental phases with testable deliverables after each sprint",
            "UML artifacts for analysis and design:",
            "    – Gantt Chart for planning and milestones",
            "    – Use Case Diagram for actor–system interactions",
            "    – Activity Diagram for key business workflows",
            "    – Sequence Diagram for controller/service collaborations",
            "Layered Client–Server SPA architecture with multi-tenant isolation",
        ],
        18,
    )

    # 8 Gantt
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "5.1 Gantt Chart")
    add_image_fit(s, gantt, Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.6))

    # 9 Use Case
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "5.2 Use Case Diagram")
    add_image_fit(s, usecase, Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.7))

    # 10 Activity
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "5.3 Activity Diagram")
    add_image_fit(s, activity, Inches(2.5), Inches(1.15), Inches(8.3), Inches(5.7))

    # 11 Sequence
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "5.4 Sequence Diagram")
    add_image_fit(s, sequence, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.7))

    # 12 Architecture
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "5.5 Architecture Overview")
    add_image_fit(s, arch, Inches(1.5), Inches(1.2), Inches(10.3), Inches(5.6))

    # 13 ERD
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "5.6 Data Model Overview (ERD)")
    add_image_fit(s, erd, Inches(0.3), Inches(1.15), Inches(12.7), Inches(5.7))

    # 14 Tools
    s = blank_slide(prs)
    add_bg(s)
    add_top_bar(s, "6. Tools Used")
    add_card(
        s,
        Inches(0.5),
        Inches(1.3),
        Inches(4.0),
        Inches(2.5),
        "Backend",
        ["Laravel (PHP)", "Eloquent ORM", "MySQL", "REST API + Sanctum"],
    )
    add_card(
        s,
        Inches(4.7),
        Inches(1.3),
        Inches(4.0),
        Inches(2.5),
        "Frontend",
        ["React (Vite)", "Material UI", "React Router", "i18n (EN / AR)"],
    )
    add_card(
        s,
        Inches(8.9),
        Inches(1.3),
        Inches(3.9),
        Inches(2.5),
        "AI & Algorithms",
        ["Google Gemini API", "Genetic Algorithm", "Constraint fitness"],
    )
    add_card(
        s,
        Inches(0.5),
        Inches(4.1),
        Inches(6.1),
        Inches(2.3),
        "Engineering Process",
        ["Spec Kit (spec-driven development)", "PHPUnit feature/unit tests", "Git / GitHub"],
    )
    add_card(
        s,
        Inches(6.9),
        Inches(4.1),
        Inches(5.9),
        Inches(2.3),
        "Design & Planning",
        ["PlantUML / Mermaid diagrams", "GanttProject", "Chapter 5 design docs"],
    )

    # 15 Demo empty
    s = blank_slide(prs)
    add_bg(s, NAVY)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, "7. Live Demo", 48, True, WHITE)

    # 16 Thank you
    s = blank_slide(prs)
    add_bg(s, NAVY)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, "Thank You", 52, True, WHITE)

    sub = s.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.7), Inches(2.0))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, "Questions & Discussion", 24, False, ACCENT)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(18)
    r2 = p2.add_run()
    set_run(r2, "Baraa Abdullah  ·  Sohayb Sabboura  ·  Obada Kahwajy", 16, False, WHITE)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    set_run(r3, "Supervisor: Anas Abdulaziz", 16, False, WHITE)

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        if i in (0, total - 2, total - 1):
            continue
        add_footer(slide, i + 1, total)

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {total}")


if __name__ == "__main__":
    build()
